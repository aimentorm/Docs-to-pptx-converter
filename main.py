# This is a sample Python script.

# Press ⌃R to execute it or replace it with your code.
# Press Double ⇧ to search everywhere for classes, files, tool windows, actions, and settings.


def print_hi(name):
    # Use a breakpoint in the code line below to debug your script.
    print(f'Hi, {name}')  # Press ⌘F8 to toggle the breakpoint.
import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches


def convert_docx_to_pptx(docx_file, pptx_file):
    # Load the .docx document
    doc = Document(docx_file)
    presentation = Presentation()

    # Iterate through paragraphs
    for para in doc.paragraphs:
        # Check if paragraph is a heading
        if para.style.name.startswith('Heading'):
            slide_layout = presentation.slide_layouts[1]  # Title and Content
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = para.text
        else:
            # If no slide exists, create one
            if len(presentation.slides) == 0:
                slide_layout = presentation.slide_layouts[1]
                slide = presentation.slides.add_slide(slide_layout)

            # Add content to the slide
            body = slide.placeholders[1].text_frame
            paragraph = body.add_paragraph()
            paragraph.text = para.text

    # Save the .pptx file
    presentation.save(pptx_file)


def main():
    docx_file = input("Enter the path to the .docx file: ")
    pptx_file = input("Enter the desired output .pptx file name: ")

    if not os.path.exists(docx_file):
        print(".docx file not found!")
        return

    convert_docx_to_pptx(docx_file, pptx_file)
    print(f"Presentation saved as {pptx_file}")
    print("Guys We have created a Document to Presentation Converter")
    print("Now We Will Create or Upload the Github or Git Repository So that our project work will be saved for reference")


# Entry point
if __name__ == "__main__":
    main()

# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    print_hi('Pycharm user, Well Done, Try with new file')


# See PyCharm help at https://www.jetbrains.com/help/pycharm/
